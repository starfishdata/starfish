import os
from typing import Dict, Any
from starfish.data_ingest.parsers.base_parser import BaseParser


class PPTParser(BaseParser):
    """Parser for PowerPoint presentations"""

    def __init__(self):
        super().__init__()
        self.supported_extensions = [".pptx"]
        self.metadata = {}

    def parse(self, file_path: str) -> str:
        """Parse a PPTX file into plain text

        Args:
            file_path: Path to the PPTX file

        Returns:
            Extracted text from the presentation
        """
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx is required for PPTX parsing. Install it with: pip install python-pptx")

        prs = Presentation(file_path)

        # Extract metadata
        self.metadata = {
            "title": prs.core_properties.title,
            "author": prs.core_properties.author,
            "created": prs.core_properties.created,
            "modified": prs.core_properties.modified,
            "slides": len(prs.slides),
        }

        # Extract text from slides
        all_text = []

        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"--- Slide {i+1} ---")

            # Get slide title
            if slide.shapes.title and slide.shapes.title.text:
                slide_text.append(f"Title: {slide.shapes.title.text}")

            # Get text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)

            all_text.append("\n".join(slide_text))

        return "\n\n".join(all_text)

    def get_metadata(self) -> Dict[str, Any]:
        """Get presentation metadata

        Returns:
            Dictionary containing presentation metadata
        """
        return self.metadata

    def is_supported(self, file_path: str) -> bool:
        """Check if the file is supported by this parser

        Args:
            file_path: Path to the file

        Returns:
            True if the file is supported, False otherwise
        """
        return os.path.splitext(file_path)[1].lower() in self.supported_extensions
